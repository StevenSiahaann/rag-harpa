import json
import pdfplumber
from docx import Document
from pptx import Presentation
from pdf2image import convert_from_path
import pytesseract
import re
from PIL import Image
import cv2
import numpy as np
import os
from langchain_community.document_loaders import AsyncChromiumLoader
from langchain_community.document_transformers import BeautifulSoupTransformer
def extract_url_from_curl(curl_command):
    """
    Extract the URL from a curl command string.
    """
    try:
        parts = curl_command.split()
        for i, part in enumerate(parts):
            if part.lower() == '-x' or part.lower() == '--url':
                return parts[i + 1]
            elif part.startswith('http'):
                return part
    except Exception as e:
        print(f"Error parsing curl command: {e}")
    return None
def get_page(urls):
    loader = AsyncChromiumLoader(urls)
    html = loader.load()

    if not html or not isinstance(html, list):
        raise ValueError("Failed to load HTML content from the URLs.")

    bs_transformer = BeautifulSoupTransformer()
    docs_transformed = bs_transformer.transform_documents(
        html, tags_to_extract=["p"], remove_unwanted_tags=["a"]
    )

    return docs_transformed

def extract_text_from_pdf(pdf_path):
    """Extract text from a PDF file with layout preservation and OCR fallback for scanned PDFs."""
    text = ""
    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text(layout=True)  # Preserve layout
            if page_text:  
                text += page_text + "\n"
                images = convert_from_path(pdf_path)
                for image in images:
                    ocr_text = pytesseract.image_to_string(image)
                    text += ocr_text + "\n"
    print("EXTRACTED TEXT: \n", clean_text(text))
    return clean_text(text)

def extract_text_from_ppt(ppt_path):
    """Extract text from a PPT file, including text in grouped shapes and tables."""
    prs = Presentation(ppt_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " ".join([cell.text for cell in row.cells])
                    text += row_text + "\n"
            if hasattr(shape, "shapes"):
                for sub_shape in shape.shapes:
                    if hasattr(sub_shape, "text"):
                        text += sub_shape.text + "\n"
    print("EXTRACTED TEXT: \n", clean_text(text))
    return clean_text(text)

def extract_text_from_docx(file_path):
    """Extract text from a .docx file."""
    doc = Document(file_path)
    text = '\n'.join([para.text for para in doc.paragraphs])
    print("EXTRACTED TEXT: \n", clean_text(text))
    return clean_text(text)

def preprocess_image(image):
    """Preprocess the image for better OCR accuracy: grayscale, contrast, and resize."""
    gray = cv2.cvtColor(np.array(image), cv2.COLOR_RGB2GRAY)
    # Resize to improve OCR accuracy if image is too small
    resized = cv2.resize(gray, None, fx=2, fy=2, interpolation=cv2.INTER_CUBIC)
    # Increase contrast
    contrast = cv2.convertScaleAbs(resized, alpha=2.0, beta=0)
    return contrast

def extract_text_from_image(image_path):
    """Extract text from an image using OCR with preprocessing."""
    image = Image.open(image_path)
    preprocessed_image = preprocess_image(image)  # Apply preprocessing for better accuracy
    text = pytesseract.image_to_string(preprocessed_image, config='--psm 6')  # Use custom OCR config for block text
    print("EXTRACTED TEXT: \n", clean_text(text))
    return clean_text(text)

def extract_text_from_txt(file_path):
    """Extract text from a text file."""
    with open(file_path, 'r', encoding='utf-8') as file:
        return file.read()

def extract_text_from_file(file_path):
    """Extract text based on the file type."""
    _, extension = os.path.splitext(file_path)
    
    if extension.lower() == ".pdf":
        return extract_text_from_pdf(file_path)
    elif extension.lower() in [".ppt", ".pptx"]:
        return extract_text_from_ppt(file_path)
    elif extension.lower() in [".jpg", ".jpeg", ".png", ".bmp", ".tiff"]:
        return extract_text_from_image(file_path)
    elif extension.lower() == ".txt":
        return extract_text_from_txt(file_path)
    elif extension.lower() == ".docx":
        return extract_text_from_docx(file_path)
    elif extension.lower() == ".json":
        extracted_text = extract_text_JSON_from_file(file_path)
        return clean_text(extracted_text)
    else:
        return None
def get_page(urls):
    loader = AsyncChromiumLoader(urls)
    html = loader.load()

    bs_transformer = BeautifulSoupTransformer()
    docs_transformed = bs_transformer.transform_documents(html, tags_to_extract=["p"], remove_unwanted_tags=["a"])

    return docs_transformed

def extract_text_JSON(json_object, indent_level=0):
    """Extract text from a JSON object with indentation."""
    lines = []
    indent = '    ' * indent_level  # Create indentation for better readability

    if isinstance(json_object, dict):
        for key, value in json_object.items():
            if isinstance(value, (dict, list)):
                nested_value = extract_text_JSON(value, indent_level + 1)
                lines.append(f'{indent}{key}: {nested_value.strip()}')  # Append key with nested values
            else:
                lines.append(f'{indent}{key}: {value}')  # Append the key-value pair directly
    elif isinstance(json_object, list):
        for index, item in enumerate(json_object):
            if isinstance(item, (dict, list)):
                nested_value = extract_text_JSON(item, indent_level + 1)
                lines.append(f'{indent}Item {index + 1}: {nested_value.strip()}')  # Append with item index
            else:
                lines.append(f'{indent}Item {index + 1}: {item}')  # Append the item directly

    return '\n'.join(lines) + '.'  # Ensure there's a period at the end of the string

def extract_text_JSON_from_file(file_path):
    """Extract text from a JSON file."""
    with open(file_path, 'r', encoding='utf-8') as json_file:
        json_object = json.load(json_file)  # This will load the whole JSON object into memory
        return extract_text_JSON(json_object)  # Call the extraction function


def clean_text(text):
    """Remove unnecessary empty lines and extra spaces from extracted text."""
    cleaned_text = re.sub(r'\s+', ' ', text.strip())
    cleaned_text = re.sub(r'\n\s*\n', '\n', cleaned_text)
    return cleaned_text

def search_json_for_keys(data, keys):
    """
    Recursively search for the given keys in a nested JSON structure, even if nested inside strings.
    :param data: The JSON data (could be dict, list, or JSON string)
    :param keys: The list of keys to search for
    :return: A dictionary with found key-value pairs
    """
    found = {}

    # If the data is a string, attempt to parse it as JSON, and continue if valid
    if isinstance(data, str):
        try:
            # Try to load the JSON string
            nested_data = json.loads(data)
            # Recursively search in this parsed structure
            return search_json_for_keys(nested_data, keys)
        except json.JSONDecodeError:
            pass  # If it's not a valid JSON string, continue

    # If the data is a dictionary, iterate over its items
    if isinstance(data, dict):
        for key, value in data.items():
            if key in keys:
                found[key] = value  # Found the key, add it to the result
            elif isinstance(value, (dict, list, str)):
                # Recursively search in nested structures, including strings that may contain JSON
                found.update(search_json_for_keys(value, keys))
    
    # If the data is a list, iterate over the items
    elif isinstance(data, list):
        for item in data:
            found.update(search_json_for_keys(item, keys))
    
    
    return found
